import os
import json
import logging
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR


def setup_logging():
    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


def create_output_folder(output_path):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)


def load_csv(input_csv):
    try:
        df = pd.read_csv(input_csv, sep="|")
        logging.info(f"Loaded CSV file with {len(df)} questions.")
        return df
    except Exception as e:
        logging.error(f"Failed to load CSV: {e}")
        raise


def create_presentation(config):
    prs = Presentation()
    prs.slide_width = Inches(config["SLIDE_WIDTH"])
    prs.slide_height = Inches(config["SLIDE_HEIGHT"])
    return prs

def add_background(slide, config):
    image_path = config["BACKGROUND_IMAGE"]
    slide_width = Inches(config["SLIDE_WIDTH"])
    slide_height = Inches(config["SLIDE_HEIGHT"])
    slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)


def add_question(slide, config, question_text):
    question_box = slide.shapes.add_textbox(Inches(config["LEFT_MARGIN"]),
                                            Inches(config["TOP_QUESTION_START"]),
                                            Inches(config["TEXT_BOX_LENGTH"]), Inches(2.0))
    fill = question_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*config["QUESTION_BG_COLOR"])

    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    question_frame.auto_size = MSO_AUTO_SIZE.NONE
    question_frame.text = question_text
    question_frame.paragraphs[0].font.size = Pt(config["FONT_SIZE_QUESTION"])
    question_frame.paragraphs[0].font.color.rgb = RGBColor(*config["QUESTION_TEXT_COLOR"])


def add_options(slide, config, options):
    option_top = Inches(config["TOP_OPTION_START"])
    option_left = Inches(config["LEFT_MARGIN"])
    row_items = 2
    counter = 0

    for key, value in options.items():
        option_box = slide.shapes.add_textbox(option_left, option_top, Inches(6.0), Inches(1.5))

        fill = option_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*config["OPTION_BG_COLOR"])

        option_frame = option_box.text_frame
        option_frame.word_wrap = True
        option_frame.text = f"{value}"
        option_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        option_frame.paragraphs[0].font.size = Pt(config["FONT_SIZE_OPTION"])
        option_frame.paragraphs[0].font.color.rgb = RGBColor(*config["OPTION_TEXT_COLOR"])

        counter += 1
        if counter % row_items == 0:
            option_top += Inches(config["LINE_SPACING"])
            option_left = Inches(config["LEFT_MARGIN"])
        else:
            option_left += Inches(config["COLUMN_SPACING"])


def add_answer(slide, config, correct_answer):
    answer_box = slide.shapes.add_textbox(Inches(config["LEFT_MARGIN"]),
                                          Inches(config["TOP_ANSWER_START"]),
                                          Inches(config["TEXT_BOX_LENGTH"]), Inches(1.0))
    fill = answer_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*config["ANSWER_BG_COLOR"])

    answer_frame = answer_box.text_frame
    answer_frame.word_wrap = True
    answer_frame.text = f"Answer: {correct_answer}"
    answer_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    answer_frame.paragraphs[0].font.size = Pt(config["FONT_SIZE_ANSWER"])
    answer_frame.paragraphs[0].font.color.rgb = RGBColor(*config["ANSWER_TEXT_COLOR"])


def save_presentation(prs, output_path):
    try:
        prs.save(output_path)
        logging.info(f"PPTX saved successfully: {output_path}")
    except Exception as e:
        logging.error(f"Failed to save PPTX: {e}")
        raise


if __name__ == "__main__":
    # Setup
    setup_logging()

    # Load config
    with open("config_generate_pptx.json", "r") as config_file:
        config = json.load(config_file)

    # Ensure output folder exists
    create_output_folder(config["OUTPUT_PPTX"])

    # Load CSV
    df = load_csv(config["INPUT_CSV"])

    # Create presentation
    prs = create_presentation(config)

    # Slide creation
    for idx, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        question_text = f"Question {idx + 1}: {row['Question']}"
        options = {
            "A": row["Option A"],
            "B": row["Option B"],
            "C": row["Option C"],
            "D": row["Option D"],
        }
        correct_answer = row["Correct Option"].strip()

        add_background(slide, config)
        add_question(slide, config, question_text)
        add_options(slide, config, options)
        add_answer(slide, config, correct_answer)

        logging.info(f"Slide {idx + 1} created.")

    save_presentation(prs, config["OUTPUT_PPTX"])
