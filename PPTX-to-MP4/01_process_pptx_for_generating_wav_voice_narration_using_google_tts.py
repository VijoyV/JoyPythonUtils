import os
import json
import logging
from io import BytesIO
from gtts import gTTS
from pptx import Presentation
from pydub import AudioSegment
from concurrent.futures import ThreadPoolExecutor

# ---------- Configuration ---------- #
DEFAULT_SILENCE_MS = 2000
OUTPUT_FORMAT = "wav"
LOG_LEVEL = logging.INFO

# ---------- Logging Setup ---------- #
logging.basicConfig(
    level=LOG_LEVEL,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

# ---------- Load Config ---------- #
def load_config(config_file):
    try:
        with open(config_file, 'r', encoding='utf-8') as file:
            config = json.load(file)
            required_keys = ['pptx_input_path_v1', 'slide_wav_narration_dir']
            for key in required_keys:
                if key not in config:
                    raise KeyError(f"Missing required config key: {key}")
            return config
    except Exception as e:
        logging.error(f"Failed to load config: {e}")
        return None

# ---------- TTS Audio Generator ---------- #
def text_to_wav_audiosegment(text):
    if not text:
        return None
    try:
        tts = gTTS(text=text, lang='en-IN')
        mp3_fp = BytesIO()
        tts.write_to_fp(mp3_fp)
        mp3_fp.seek(0)
        audio = AudioSegment.from_file(mp3_fp, format="mp3")
        return audio.set_frame_rate(44100).set_channels(2)
    except Exception as e:
        logging.error(f"Error during TTS for text: '{text[:30]}...': {e}")
        return None

# ---------- Slide Processor ---------- #
def process_slide(slide_no, slide, output_dir, silence_duration):
    try:
        logging.info(f"Processing Slide {slide_no + 1}")
        text_shapes = [
            shape for shape in slide.shapes
            if hasattr(shape, 'text') and shape.name
        ]
        sorted_shapes = sorted(text_shapes, key=lambda s: s.name)

        segments = []
        for shape in sorted_shapes:
            text = shape.text.strip()
            if text:
                logging.debug(f"  Shape: {shape.name} | Text: {text}")
                audio = text_to_wav_audiosegment(text)
                if audio:
                    segments.append(audio + AudioSegment.silent(duration=silence_duration))

        if segments:
            combined = sum(segments[1:], segments[0])
            out_path = os.path.join(output_dir, f"slide_{slide_no + 1}_narration.{OUTPUT_FORMAT}")
            combined.export(out_path, format=OUTPUT_FORMAT)
            logging.info(f"  Saved narration: {out_path}")
        else:
            logging.warning(f"  No valid text found in Slide {slide_no + 1}")
    except Exception as e:
        logging.error(f"Error processing slide {slide_no + 1}: {e}")

# ---------- Presentation Processor ---------- #
def process_presentation_for_tts(config):
    pptx_path = config['pptx_input_path_v1']
    output_dir = config['slide_wav_narration_dir']
    silence_duration = config.get('silence_duration_ms', DEFAULT_SILENCE_MS)
    parallel = config.get('parallel_processing', False)

    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)

    if parallel:
        with ThreadPoolExecutor() as executor:
            futures = [
                executor.submit(process_slide, i, slide, output_dir, silence_duration)
                for i, slide in enumerate(prs.slides)
            ]
            for future in futures:
                future.result()  # wait for all
    else:
        for i, slide in enumerate(prs.slides):
            process_slide(i, slide, output_dir, silence_duration)

# ---------- Main ---------- #
if __name__ == "__main__":
    config_path = '01_config.json'
    config = load_config(config_path)
    if config:
        process_presentation_for_tts(config)
