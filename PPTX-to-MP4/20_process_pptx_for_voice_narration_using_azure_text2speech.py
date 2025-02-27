import os
import json
import requests
from pptx import Presentation

# --- Load Configuration ---
config_file = "config-azure.json"
with open(config_file, "r") as f:
    config = json.load(f)

# Azure TTS configuration
speech_key = config["azure"]["speech_key"]
service_region = config["azure"]["service_region"]
voice_selection = config["azure"]["voice"].lower()  # "male" or "female"
speech_rate = config["azure"]["rate"]

# Read the TTS endpoint from config and build the full URL.
# For Azure TTS, the correct endpoint is:
#   https://<region>.tts.speech.microsoft.com/cognitiveservices/v1
tts_endpoint = config["azure"].get("tts_endpoint", "")
tts_url = tts_endpoint.rstrip("/") + "/cognitiveservices/v1"

# PPTX and output configuration
pptx_filename = config["pptx"]["input_file"]
output_dir = config["output"]["audio_directory"]

# Select the appropriate voice name for English Indian accent.
if voice_selection == "female":
    voice_name = "en-IN-NeerjaNeural"
else:
    voice_name = "en-IN-PrabhatNeural"

# Create the output directory if it does not exist.
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# --- Set up the REST API endpoint and headers ---
headers = {
    "Ocp-Apim-Subscription-Key": speech_key,
    "Content-Type": "application/ssml+xml",
    "X-Microsoft-OutputFormat": "riff-24khz-16bit-mono-pcm",
    "User-Agent": "PythonTTSClient"
}

def synthesize_text(text, voice_name, rate, output_path):
    # Construct SSML as a single line with proper namespace.
    ssml = (
        f'<speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis" xml:lang="en-IN">'
        f'<voice name="{voice_name}">'
        f'<prosody rate="{rate}">{text}</prosody>'
        f'</voice></speak>'
    )
    response = requests.post(tts_url, headers=headers, data=ssml.encode("utf-8"))
    if response.status_code == 200:
        with open(output_path, "wb") as audio_file:
            audio_file.write(response.content)
        print(f"Audio saved: {output_path}")
    else:
        print(f"Error synthesizing audio. Status code: {response.status_code}")
        print(response.text)

# --- Extract Text from Each Slide using python-pptx ---
prs = Presentation(pptx_filename)
slide_texts = []
for idx, slide in enumerate(prs.slides, start=1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    slide_text = "\n".join(texts)
    slide_texts.append(slide_text)
    print(f"Extracted text from slide {idx}.")

# --- Generate Audio for Each Slide via REST API ---
for idx, text in enumerate(slide_texts, start=1):
    if not text:
        print(f"Slide {idx} is empty; skipping audio generation.")
        continue

    audio_filename = os.path.join(output_dir, f"slide_{idx}.wav")
    print(f"Generating audio for slide {idx} with speech rate {speech_rate}...")
    synthesize_text(text, voice_name, speech_rate, audio_filename)
