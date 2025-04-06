import os
from pptx import Presentation
from config_loader import load_config

def extract_and_save_text_from_pptx(pptx_file, output_text_dir):
    """Extract text from PPTX slides, sorting text-containing objects by name, and save to text files."""
    os.makedirs(output_text_dir, exist_ok=True)

    prs = Presentation(pptx_file)

    for idx, slide in enumerate(prs.slides, start=1):
        # Collect text shapes, sorting by their shape name
        text_shapes = sorted(
            [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()],
            key=lambda shape: shape.name if hasattr(shape, "name") else ""
        )

        slide_text = "\n".join([shape.text.strip() for shape in text_shapes])
        text_filename = os.path.join(output_text_dir, f"slide_{idx}.txt")

        with open(text_filename, "w", encoding="utf-8") as text_file:
            text_file.write(slide_text)

        print(f"Extracted and saved text for slide {idx}: {text_filename}")

def read_text_from_files(text_dir):
    """Read slide texts from stored text files instead of extracting from PPTX."""
    slide_texts = []
    text_files = sorted(
        [f for f in os.listdir(text_dir) if f.startswith("slide_") and f.endswith(".txt")],
        key=lambda x: int(x.split("_")[1].split(".")[0])  # Sort by slide number
    )

    for text_file in text_files:
        with open(os.path.join(text_dir, text_file), "r", encoding="utf-8") as f:
            slide_texts.append(f.read().strip())

    return slide_texts
