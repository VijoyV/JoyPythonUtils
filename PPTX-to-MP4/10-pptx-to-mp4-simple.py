import json
import os
from pptx import Presentation
from moviepy.editor import *

# Function to create video from PPTX and audio
def create_video(pptx_path, audio_config_path, output_path):
    # Load the PPTX
    prs = Presentation(pptx_path)

    # Load audio config
    with open(audio_config_path, 'r') as f:
        audio_config = json.load(f)

    slides_clips = []

    # Iterate through slides
    for i, slide in enumerate(prs.slides):
        slide_img = f"slide_{i+1}.png"

        # Save slide as image
        slide.shapes._spTree.remove(slide.shapes.title._element) if slide.shapes.title else None
        prs.save(slide_img)

        # Get corresponding audio
        audio_file = audio_config.get(f"slide_{i+1}")
        if not audio_file or not os.path.exists(audio_file):
            print(f"Audio for slide {i+1} not found. Skipping.")
            continue

        # Load audio and create image clip
        audio_clip = AudioFileClip(audio_file)
        img_clip = ImageClip(slide_img).set_duration(audio_clip.duration).set_audio(audio_clip)

        slides_clips.append(img_clip)

    # Concatenate all clips
    final_clip = concatenate_videoclips(slides_clips, method="compose")

    # Write to MP4
    final_clip.write_videofile(output_path, fps=24)

    # Cleanup images
    for i in range(len(slides_clips)):
        os.remove(f"slide_{i+1}.png")

# Example usage
create_video("presentation.pptx", "audio_config.json", "output_video.mp4")

"""

Config File:

{
  "slide_1": "audio/slide1_audio.wav",
  "slide_2": "audio/slide2_audio.wav",
  "slide_3": "audio/slide3_audio.wav",
  "slide_4": "audio/slide4_audio.wav",
  "slide_5": "audio/slide5_audio.wav"
}

"""