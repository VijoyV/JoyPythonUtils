from pptx import Presentation
from pptx.util import Inches
import glob

# Path to the directory containing .jpg files
jpg_dir = 'C:\\Learning\\Google-IntroToGenAI\\images'

# Create a new PowerPoint presentation
prs = Presentation()

# Get a list of .jpg files in the directory
jpg_files = glob.glob(f"{jpg_dir}/*.jpg")

# Add each .jpg file as a slide in the presentation
for jpg_file in jpg_files:
    slide_layout = prs.slide_layouts[5]  # Use layout 5 (Title and Content) for each slide
    slide = prs.slides.add_slide(slide_layout)
    content = slide.shapes.add_picture(jpg_file, Inches(0), Inches(0), width=Inches(10))  # Adjust width as needed

# Save the presentation
pptx_file = 'C:\\Learning\\Google-IntroToGenAI\\IntroToGoogleGenAI.pptx'
prs.save(pptx_file)

print(f"PowerPoint presentation saved as {pptx_file}")
