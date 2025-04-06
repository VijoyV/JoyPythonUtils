import os
import logging
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE  # Import auto-size options

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Constants
INPUT_CSV = "./input_csv/Judges_Chapter_11_QnA.csv"
OUTPUT_PPTX = "./output_pptx/Judges_Chapter_11_MCQs.pptx"
FONT_NAME = "Arial"
FONT_SIZE_QUESTION = Pt(32)
FONT_SIZE_OPTION = Pt(32)
FONT_SIZE_ANSWER = Pt(28)
SLIDE_WIDTH = Inches(13.33)  # Widescreen 16:9
SLIDE_HEIGHT = Inches(7.5)

# Ensure output folder exists
os.makedirs(os.path.dirname(OUTPUT_PPTX), exist_ok=True)

# Load CSV
try:
    df = pd.read_csv(INPUT_CSV, sep="|")
    logging.info(f"Loaded CSV file with {len(df)} questions.")
except Exception as e:
    logging.error(f"Failed to load CSV: {e}")
    raise

# Create widescreen presentation
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Slide creation
for idx, row in df.iterrows():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

    # Get question and options
    question_text = f"Question {idx+1}: {row['Question']}"
    options = {
        "A": row["Option A"],
        "B": row["Option B"],
        "C": row["Option C"],
        "D": row["Option D"]
    }
    correct_answer = row["Correct Option"].strip()

    # Position variables
    left_margin = Inches(1)
    top_question = Inches(1)
    top_option_start = Inches(3)
    line_spacing = Inches(0.75)

    # Add Question Box with Wrapping
    question_box = slide.shapes.add_textbox(left_margin, Inches(0.5), Inches(12), Inches(2.4))  # taller to allow wrapping
    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    question_frame.auto_size = MSO_AUTO_SIZE.NONE
    question_frame.text = question_text
    question_frame.paragraphs[0].font.size = FONT_SIZE_QUESTION


    # Add Options
    option_top = top_option_start
    for key, value in options.items():
        option_box = slide.shapes.add_textbox(left_margin, option_top, Inches(12), Inches(0.75))
        option_frame = option_box.text_frame
        option_frame.text = f"{value}"
        option_frame.paragraphs[0].font.size = FONT_SIZE_OPTION
        option_top += line_spacing

    # Add Answer Box
    answer_box = slide.shapes.add_textbox(left_margin, option_top + Inches(0.5), Inches(12), Inches(0.75))
    answer_frame = answer_box.text_frame
    answer_frame.word_wrap = True
    answer_frame.text = f"Answer: {correct_answer}"
    answer_frame.paragraphs[0].font.size = FONT_SIZE_ANSWER
    answer_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # Red color for answer

    logging.info(f"Slide {idx+1} created.")

# Save the presentation
try:
    prs.save(OUTPUT_PPTX)
    logging.info(f"PPTX saved successfully: {OUTPUT_PPTX}")
except Exception as e:
    logging.error(f"Failed to save PPTX: {e}")
    raise
