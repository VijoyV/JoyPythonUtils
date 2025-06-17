import json
import logging
import os

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


def setup_logging():
    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


def create_output_folder(output_path):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)


def load_csv(input_csv):
    try:
        df = pd.read_csv(input_csv, sep="|")
        logging.info(f"Loaded CSV file with {len(df)} questions.")
        return df
    except Exception as e:
        logging.error(f"Failed to load CSV: {e}")
        raise


def create_presentation(config):
    prs = Presentation()
    prs.slide_width = Inches(config["SLIDE_WIDTH"])
    prs.slide_height = Inches(config["SLIDE_HEIGHT"])
    return prs


def add_background(slide, config):
    image_path = config["BACKGROUND_IMAGE"]
    slide_width = Inches(config["SLIDE_WIDTH"])
    slide_height = Inches(config["SLIDE_HEIGHT"])
    slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)


def get_color(color_value):
    """Supports named colors, hex codes (#RRGGBB), and RGB lists [R,G,B]"""
    color_map = {
        # Basic Colors
        'black': RGBColor(0, 0, 0),
        'white': RGBColor(255, 255, 255),
        'red': RGBColor(255, 0, 0),
        'green': RGBColor(0, 128, 0),
        'blue': RGBColor(0, 0, 255),
        'yellow': RGBColor(255, 255, 0),
        'cyan': RGBColor(0, 255, 255),
        'magenta': RGBColor(255, 0, 255),

        # Blues
        'lightblue': RGBColor(173, 216, 230),
        'skyblue': RGBColor(135, 206, 235),
        'royalblue': RGBColor(65, 105, 225),
        'darkblue': RGBColor(0, 0, 139),
        'navy': RGBColor(0, 0, 128),

        # Greens
        'lime': RGBColor(0, 255, 0),
        'forestgreen': RGBColor(34, 139, 34),
        'seagreen': RGBColor(46, 139, 87),
        'olive': RGBColor(128, 128, 0),

        # Reds/Pinks
        'crimson': RGBColor(220, 20, 60),
        'pink': RGBColor(255, 192, 203),
        'hotpink': RGBColor(255, 105, 180),

        # Purples
        'purple': RGBColor(128, 0, 128),
        'lavender': RGBColor(230, 230, 250),
        'violet': RGBColor(238, 130, 238),

        # Oranges/Browns
        'orange': RGBColor(255, 165, 0),
        'gold': RGBColor(255, 215, 0),
        'brown': RGBColor(165, 42, 42),
        'tan': RGBColor(210, 180, 140),

        # Grays
        'silver': RGBColor(192, 192, 192),
        'gray': RGBColor(128, 128, 128),
        'darkgray': RGBColor(169, 169, 169)
    }

    # Handle hex codes (#RRGGBB)
    if isinstance(color_value, str) and color_value.startswith('#'):
        hex = color_value.lstrip('#')
        return RGBColor(int(hex[0:2], 16), int(hex[2:4], 16), int(hex[4:6], 16))

    # Handle RGB lists
    elif isinstance(color_value, list):
        return RGBColor(*color_value)

    # Handle named colors (case-insensitive)
    return color_map.get(color_value.lower(), RGBColor(0, 0, 0))  # Default: black


def add_question(slide, config, question_text):
    question_box = slide.shapes.add_textbox(Inches(config["LEFT_MARGIN"]),
                                            Inches(config["TOP_QUESTION_START"]),
                                            Inches(config["TEXT_BOX_LENGTH"]), Inches(2.0))

    if config["USE_BACKGROUND_FILL"].get("question", False):
        fill = question_box.fill
        fill.solid()
        fill.fore_color.rgb = get_color(config["QUESTION_BG_COLOR"])
    else:
        question_box.fill.background()  # Transparent background

    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    question_frame.auto_size = MSO_AUTO_SIZE.NONE
    question_frame.text = question_text
    question_frame.paragraphs[0].font.name = config["FONT_NAME"]
    question_frame.paragraphs[0].font.size = Pt(config["FONT_SIZE_QUESTION"])
    question_frame.paragraphs[0].font.color.rgb = get_color(config["QUESTION_TEXT_COLOR"])


def add_options(slide, config, options):
    option_top = Inches(config["TOP_OPTION_START"])
    option_left = Inches(config["LEFT_MARGIN"])
    row_items = 2
    counter = 0

    for key, value in options.items():
        option_box = slide.shapes.add_textbox(option_left, option_top, Inches(5.75), Inches(1.25))

        if config["USE_BACKGROUND_FILL"].get("options", False):
            fill = option_box.fill
            fill.solid()
            fill.fore_color.rgb = get_color(config["OPTION_BG_COLOR"])
        else:
            option_box.fill.background()  # Transparent background

        option_frame = option_box.text_frame
        option_frame.word_wrap = True
        option_frame.text = f"{value}"
        option_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        option_frame.paragraphs[0].font.name = config["FONT_NAME"]
        option_frame.paragraphs[0].font.size = Pt(config["FONT_SIZE_OPTION"])
        option_frame.paragraphs[0].font.color.rgb = get_color(config["OPTION_TEXT_COLOR"])

        counter += 1
        if counter % row_items == 0:
            option_top += Inches(config["LINE_SPACING"])
            option_left = Inches(config["LEFT_MARGIN"])
        else:
            option_left += Inches(config["COLUMN_SPACING"])


def add_answer(slide, config, correct_answer):
    answer_box = slide.shapes.add_textbox(Inches(config["LEFT_MARGIN"]),
                                          Inches(config["TOP_ANSWER_START"]),
                                          Inches(config["TEXT_BOX_LENGTH"]), Inches(1.0))

    if config["USE_BACKGROUND_FILL"].get("answer", False):
        fill = answer_box.fill
        fill.solid()
        fill.fore_color.rgb = get_color(config["ANSWER_BG_COLOR"])
    else:
        answer_box.fill.background()  # Transparent background

    answer_frame = answer_box.text_frame
    answer_frame.word_wrap = True
    answer_frame.text = f"Answer: {correct_answer}"
    answer_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    answer_frame.paragraphs[0].font.name = config["FONT_NAME"]
    answer_frame.paragraphs[0].font.size = Pt(config["FONT_SIZE_ANSWER"])
    answer_frame.paragraphs[0].font.color.rgb = get_color(config["ANSWER_TEXT_COLOR"])


def save_presentation(prs, output_path):
    try:
        prs.save(output_path)
        logging.info(f"PPTX saved successfully: {output_path}")
    except Exception as e:
        logging.error(f"Failed to save PPTX: {e}")
        raise


if __name__ == "__main__":
    # Setup
    setup_logging()

    # Load config
    with open("config_generate_pptx.json", "r") as config_file:
        config = json.load(config_file)

    # Ensure output folder exists
    create_output_folder(config["OUTPUT_PPTX"])

    # Load CSV
    df = load_csv(config["INPUT_CSV"])

    # Create presentation
    prs = create_presentation(config)

    # Slide creation
    for idx, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        question_text = f"Question {idx + 1}: {row['Question']}"
        options = {
            "A": row["Option A"],
            "B": row["Option B"],
            "C": row["Option C"],
            "D": row["Option D"],
        }
        correct_answer = row["Correct Option"].strip()

        add_background(slide, config)
        add_question(slide, config, question_text)
        add_options(slide, config, options)
        add_answer(slide, config, correct_answer)

        logging.info(f"Slide {idx + 1} created.")

    save_presentation(prs, config["OUTPUT_PPTX"])
