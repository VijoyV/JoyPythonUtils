from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_background_image(slide, image_path, prs):
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)


def add_first_slide(prs, background_image_path, title_text, subtitle_text):
    first_slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    first_slide = prs.slides.add_slide(first_slide_layout)
    add_background_image(first_slide, background_image_path, prs)

    # Adding title
    title_shape = first_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    title_text_frame = title_shape.text_frame
    title_text_frame.text = title_text
    title_text_frame.paragraphs[0].font.size = Pt(56)
    title_text_frame.paragraphs[0].font.bold = True
    title_text_frame.paragraphs[0].font.color.rgb = RGBColor(20, 150, 25)
    title_text_frame.word_wrap = True  # Enable word wrap

    # Adding subtitle
    subtitle_shape = first_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    subtitle_text_frame = subtitle_shape.text_frame
    subtitle_text_frame.text = subtitle_text
    subtitle_text_frame.paragraphs[0].font.size = Pt(30)
    subtitle_text_frame.paragraphs[0].font.color.rgb = RGBColor(55, 125, 55)
    subtitle_text_frame.word_wrap = True  # Enable word wrap

    # Move the first slide to the beginning
    xml_slides = prs.slides._sldIdLst  # Access the slide ID list
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(0, slides[-1])


def add_last_slide(prs, background_image_path, last_slide_message):
    last_slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    last_slide = prs.slides.add_slide(last_slide_layout)
    add_background_image(last_slide, background_image_path, prs)

    # Adding message
    message_shape = last_slide.shapes.add_textbox(Inches(4), Inches(5), Inches(9), Inches(3))
    message_text_frame = message_shape.text_frame
    message_text_frame.text = last_slide_message
    message_text_frame.paragraphs[0].font.size = Pt(40)
    message_text_frame.paragraphs[0].font.bold = True
    message_text_frame.paragraphs[0].font.color.rgb = RGBColor(20, 150, 25)
    message_text_frame.word_wrap = True  # Enable word wrap


def main():
    # Variables (Assume these are provided)
    background_image_path_beg = 'input/BG-Beg.jpg'
    background_image_path_end = 'input/BG-End.png'
    title_text = "Logos Bible Quiz 2024"
    subtitle_text = "Judges - Chapter 03: Important Questions"
    last_slide_message = "Thank you for watching this video! Please subscribe to the channel for more videos like this."
    pptx_file_path = './output/Judges_Chapter_03_MCQ_V0.pptx'
    output_pptx_file_path = './output/Judges_Chapter_03_MCQ_V1.pptx'

    # Load the presentation
    prs = Presentation(pptx_file_path)

    # Add first slide at the beginning
    add_first_slide(prs, background_image_path_beg, title_text, subtitle_text)

    # Add last slide at the end
    add_last_slide(prs, background_image_path_end, last_slide_message)

    # Save the modified presentation

    prs.save(output_pptx_file_path)

    print(f"Presentation saved as {output_pptx_file_path}")


if __name__ == "__main__":
    main()
