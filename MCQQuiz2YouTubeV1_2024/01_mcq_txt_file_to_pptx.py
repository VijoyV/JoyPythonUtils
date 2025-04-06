from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.util import Inches


def parse_mcqs(file_path):
    mcqs = []
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            lines = file.readlines()
    except FileNotFoundError:
        print(f"Error: File not found at {file_path}")
        return mcqs
    except Exception as e:
        print(f"Error reading file: {e}")
        return mcqs

    mcq = {}
    for line in lines:
        line = line.strip()
        if line.startswith("Question"):
            if mcq:
                mcqs.append(mcq)
                mcq = {}
            mcq['question'] = line.split(':', 1)[1].strip()
        elif line.startswith("Options:"):
            mcq['options'] = []
        elif line.startswith('A)') or line.startswith('B)') or line.startswith('C)') or line.startswith('D)'):
            mcq['options'].append(line)
        elif line.startswith("Answer:"):
            mcq['answer'] = line.split(':', 1)[1].strip()
    if mcq:
        mcqs.append(mcq)
    return mcqs


def set_font_properties(paragraph, size, color=None):
    paragraph.font.size = Pt(size)
    if color:
        paragraph.font.color.rgb = color


def add_slide(prs, mcq, index, slide_title, watermark_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Add watermark image
    slide.shapes.add_picture(watermark_path, Inches(0), Inches(0),
                             width=prs.slide_width, height=prs.slide_height)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.1), Inches(10), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title = title_frame.add_paragraph()
    title.text = slide_title
    set_font_properties(title, 36, RGBColor(102, 55, 104))

    # Add question
    question_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(12), Inches(3))
    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    p = question_frame.add_paragraph()
    p.text = f"Question - {index + 1}: {mcq['question']}"
    set_font_properties(p, 24, RGBColor(0, 102, 224))

    # Add options
    options_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(12), Inches(3))
    options_frame = options_box.text_frame
    options_frame.word_wrap = True
    for option in mcq['options']:
        p = options_frame.add_paragraph()
        p.text = option
        p.level = 0
        set_font_properties(p, 24)

    # Add answer
    answer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(12), Inches(1))
    answer_frame = answer_box.text_frame
    answer_frame.word_wrap = True
    p = answer_frame.add_paragraph()
    p.text = f"ANS: {mcq['answer']}"
    set_font_properties(p, 24, RGBColor(0, 102, 204))


def generate_ppt(mcqs, output_path, slide_title, watermark_path):
    prs = Presentation()

    # Set the slide orientation to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for index, mcq in enumerate(mcqs):
        add_slide(prs, mcq, index, slide_title, watermark_path)

    prs.save(output_path)
    print(f"PowerPoint presentation saved as {output_path}")


def main():
    slide_title = "Judges : Chapter 03"
    file_path = './input/Judges_Chapter_03_MCQ.txt'
    output_path = './output/Judges_Chapter_03_MCQ_V0.pptx'
    watermark_path = 'input/BG-Normal.png'  # Path to your watermark image

    mcqs = parse_mcqs(file_path)
    generate_ppt(mcqs, output_path, slide_title, watermark_path)


if __name__ == "__main__":
    main()
