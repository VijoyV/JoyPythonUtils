import os
import time
import json
from gtts import gTTS
from pptx import Presentation
from pydub import AudioSegment

# Function to load the configuration from a JSON file
def load_config(config_file):
    """Load and return the configuration from a JSON file."""
    try:
        with open(config_file, 'r', encoding='utf-8') as file:
            return json.load(file)
    except (FileNotFoundError, json.JSONDecodeError) as e:
        print(f"Error loading config file: {e}")
        return None

def generate_tts(text, output_file):
    """Generate speech from text and save it as an MP3 file."""
    if not text:
        print(f"No text provided for {output_file}. Skipping TTS generation.")
        return
    tts = gTTS(text=text, lang='en-IN')
    tts.save(output_file)
    print(f"Generated TTS audio file: {output_file}")

def add_silence_and_combine(audio_files, silence_duration, output_file):
    """Combine multiple audio files with silences and export the result."""
    combined_audio = AudioSegment.silent(duration=0)

    for audio_file in audio_files:
        if os.path.exists(audio_file):
            audio = AudioSegment.from_mp3(audio_file)
            combined_audio += audio + AudioSegment.silent(duration=silence_duration)
        else:
            print(f"Warning: {audio_file} not found.")  # Debugging: Check for missing audio files

    combined_audio.export(output_file, format="mp3")
    print(f"Combined TTS audio file: {output_file}")

    # Optionally, remove the intermediate files
    for audio_file in audio_files:
        if os.path.exists(audio_file):
            os.remove(audio_file)

def process_slide(slide_no, slide, narration_output_dir):
    """Process a slide and generate narration MP3s with silence."""
    shape_texts = []

    # Extract text from all shapes that are textboxes
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            shape_text = shape.text.strip()
            if shape_text:
                print(f"Slide {slide_no + 1} Text: {shape_text}")  # Debugging: print the text of each shape
                shape_texts.append(shape_text)

    # Generate TTS for each text box and add to audio files list
    audio_files = []
    for idx, text in enumerate(shape_texts):
        if text:
            audio_file = os.path.join(narration_output_dir, f"slide_{slide_no + 1}_narration_part_{idx + 1}.mp3")
            generate_tts(text, audio_file)
            audio_files.append(audio_file)

    # Combine with 5-second silences between texts
    silence_duration = 5000  # 5 seconds silence
    if audio_files:
        combined_audio_file = os.path.join(narration_output_dir, f"slide_{slide_no + 1}_narration.mp3")
        add_silence_and_combine(audio_files, silence_duration, combined_audio_file)

    time.sleep(2)  # Delay for processing

def process_presentation_for_tts(config):
    """Process each slide in a PowerPoint presentation, generate TTS audio, and save it."""
    presentation_path = config['ppt_output_path_v1']
    narration_output_dir = config['slide_mp3_narration_dir']

    os.makedirs(narration_output_dir, exist_ok=True)

    prs = Presentation(presentation_path)

    for i, slide in enumerate(prs.slides):
        print(f"Processing Slide {i + 1}")
        process_slide(i, slide, narration_output_dir)

if __name__ == "__main__":
    config_file = 'config.json'
    config = load_config(config_file)
    if config:
        process_presentation_for_tts(config)
