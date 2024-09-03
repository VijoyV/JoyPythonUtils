import os
import json
from gtts import gTTS
from pptx import Presentation
from pydub import AudioSegment

# Function to load the configuration from a JSON file
def load_config(config_file):
    try:
        with open(config_file, 'r', encoding='utf-8') as file:
            config = json.load(file)
        return config
    except FileNotFoundError:
        print(f"Error: Config file not found at {config_file}")
        return None
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}")
        return None

def generate_tts(text, output_file):
    """Generate speech from text and save it as an MP3 file."""
    tts = gTTS(text=text, lang='en')
    tts.save(output_file)
    print(f"Generated TTS audio file: {output_file}")

def process_presentation_for_tts(config):
    """Process each slide in a PowerPoint presentation, generate TTS audio for shapes 3, 4, and 5 with a 3-second gap, and save it."""
    presentation_path = config['ppt_output_path_v1']
    narration_output_dir = config['slide_mp3_narration_dir']

    # Ensure the output directory exists
    if not os.path.exists(narration_output_dir):
        os.makedirs(narration_output_dir)

    prs = Presentation(presentation_path)

    for i, slide in enumerate(prs.slides):
        # Skip slides without sufficient shapes
        if len(slide.shapes) >= 5:
            print(f"Slide {i+1} OK: has {len(slide.shapes)} shapes")
        else:
            print(f"Slide {i+1} skipped: less than 5 shapes found.")
            continue

        # Extract text from shapes 3, 4, and 5
        narration_text_3_4 = ""
        shape_3 = slide.shapes[2]
        shape_4 = slide.shapes[3]
        shape_5 = slide.shapes[4]

        if hasattr(shape_3, "text"):
            narration_text_3_4 += shape_3.text + " "

        if hasattr(shape_4, "text"):
            narration_text_3_4 += shape_4.text + " "

        if hasattr(shape_5, "text"):
            shape_5_text = shape_5.text
        else:
            shape_5_text = ""

        # Generate TTS audio for shapes 3 and 4
        if narration_text_3_4:
            audio_file_3_4 = os.path.join(narration_output_dir, f"slide_{i + 1}_narration_part_1.mp3")
            generate_tts(narration_text_3_4, audio_file_3_4)
        else:
            print(f"No text found in shapes 3 and 4 on slide {i + 1}.")
            continue

        # Generate TTS audio for shape 5
        if shape_5_text:
            audio_file_5 = os.path.join(narration_output_dir, f"slide_{i + 1}_narration_part_2.mp3")
            generate_tts(shape_5_text, audio_file_5)
        else:
            print(f"No text found in shape 5 on slide {i + 1}.")
            continue

        # Create a 3-second silence
        silence_1 = AudioSegment.silent(duration=2000)  # 2 seconds
        silence_2 = AudioSegment.silent(duration=1500)  # 1.5 seconds

        # Combine the TTS audio files and silence
        narration_audio_3_4 = AudioSegment.from_mp3(audio_file_3_4)
        narration_audio_5 = AudioSegment.from_mp3(audio_file_5)
        combined_audio = narration_audio_3_4 + silence_1 + narration_audio_5 + silence_2

        # Save the combined audio
        combined_audio_file = os.path.join(narration_output_dir, f"slide_{i + 1}_narration.mp3")
        combined_audio.export(combined_audio_file, format="mp3")
        print(f"Combined TTS audio file: {combined_audio_file}")

        # Optionally, remove the intermediate files
        os.remove(audio_file_3_4)
        os.remove(audio_file_5)

if __name__ == "__main__":
    # Load configuration
    config_file = 'config.json'
    config = load_config(config_file)
    if config is None:
        exit(1)  # Exit if config could not be loaded

    # Process the presentation for TTS
    process_presentation_for_tts(config)
