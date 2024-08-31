import os
import json
from gtts import gTTS
from pptx import Presentation

# Function to load the configuration from a JSON file
def load_config(config_file):
    try:
        with open(config_file, 'r', encoding='utf-8') as file:
            config = json.load(file)
        return config
    except FileNotFoundError:
        print(f"Error: Config file not found at {config_file}")
        return None
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}")
        return None

def generate_tts(text, output_file):
    """Generate speech from text and save it as an MP3 file."""
    tts = gTTS(text=text, lang='en')
    tts.save(output_file)
    print(f"Generated TTS audio file: {output_file}")

def process_presentation_for_tts(config):
    """Process each slide in a PowerPoint presentation, generate TTS audio for shapes 3 and 4, and save it."""
    presentation_path = config['ppt_output_path_v1']
    narration_output_dir = config['slide_mp3_narration_dir']

    # Ensure the output directory exists
    if not os.path.exists(narration_output_dir):
        os.makedirs(narration_output_dir)

    prs = Presentation(presentation_path)

    for i, slide in enumerate(prs.slides):
        # Skip slides without sufficient shapes
        if len(slide.shapes) >= 5:
            print(f"Slide {i+1} OK: has {len(slide.shapes)} shapes")
        else:
            print(f"Slide {i+1} skipped: less than 5 shapes found.")
            continue

        # Extract text from shapes 2, 3 and 4
        narration_text = ""
        shapes_to_process = [2, 3, 4]
        for shape_index in shapes_to_process:
            shape = slide.shapes[shape_index]
            if hasattr(shape, "text"):
                narration_text += shape.text + " "

        # Generate TTS audio file
        if narration_text:
            audio_file = os.path.join(narration_output_dir, f"slide_{i + 1}_narration.mp3")
            generate_tts(narration_text, audio_file)
        else:
            print(f"No text found in shapes 2 and 3 on slide {i + 1}.")

if __name__ == "__main__":
    # Load configuration
    config_file = 'config.json'
    config = load_config(config_file)
    if config is None:
        exit(1)  # Exit if config could not be loaded

    # Process the presentation for TTS
    process_presentation_for_tts(config)
