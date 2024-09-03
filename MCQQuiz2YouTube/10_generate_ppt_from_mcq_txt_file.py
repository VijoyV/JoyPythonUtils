import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Function to load the configuration from a JSON file
def load_config(config_file):
    try:
        with open(config_file, 'r', encoding='utf-8') as file:
            config = json.load(file)
        return config
    except FileNotFoundError:
        print(f"Error: Config file not found at {config_file}")
        return None
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}")
        return None

# Function to parse MCQs from a text file
def parse_mcqs(file_path):
    mcqs = []
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            lines = file.readlines()
    except FileNotFoundError:
        print(f"Error: File not found at {file_path}")
        return mcqs
    except Exception as e:
        print(f"Error reading file: {e}")
        return mcqs

    mcq = {}
    for line in lines:
        line = line.strip()
        if line.startswith("Question:"):
            if mcq:
                mcqs.append(mcq)
                mcq = {}
            mcq['question'] = line.split(':', 1)[1].strip()
        elif line.startswith("Options:"):
            mcq['options'] = []
        elif line.startswith('A)') or line.startswith('B)') or line.startswith('C)') or line.startswith('D)'):
            mcq['options'].append(line)
        elif line.startswith("Answer:"):
            mcq['answer'] = line.split(':', 1)[1].strip()
    if mcq:
        mcqs.append(mcq)
    return mcqs

# Function to set font properties for text paragraphs
def set_font_properties(paragraph, size, color=None):
    paragraph.font.size = Pt(size)
    if color:
        paragraph.font.color.rgb = color

# Function to add a slide with an MCQ to the presentation
def add_slide(prs, mcq, index, slide_title, watermark_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Add watermark image
    slide.shapes.add_picture(watermark_path, Inches(0), Inches(0),
                             width=prs.slide_width, height=prs.slide_height)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.1), Inches(10), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title = title_frame.add_paragraph()
    title.text = slide_title
    set_font_properties(title, 36, RGBColor(102, 55, 104))

    # Check the shape type
    # if title_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Title Box is a TextBox.")
    # else:
    #     print(f"The shape of Title Box is not a TextBox, it is of type {title_box.shape_type}")


    # Add question
    question_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(12), Inches(3))
    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    p = question_frame.add_paragraph()
    p.text = f"Question - {index + 1}: {mcq['question']}"
    set_font_properties(p, 24, RGBColor(0, 102, 224))

    # Check the shape type
    # if question_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Question Box is a TextBox.")
    # else:
    #     print(f"The shape of Question Box is not a TextBox, it is of type {question_box.shape_type}")


    # Add options
    options_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(12), Inches(3))
    options_frame = options_box.text_frame
    options_frame.word_wrap = True
    for option in mcq['options']:
        p = options_frame.add_paragraph()
        p.text = option
        p.level = 0
        set_font_properties(p, 24)

    # Check the shape type
    # if options_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Options Box is a TextBox.")
    # else:
    #     print(f"The shape of Options Box is not a TextBox, it is of type {options_box.shape_type}")

    # Add answer
    answer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(12), Inches(1))
    answer_frame = answer_box.text_frame
    answer_frame.word_wrap = True
    p = answer_frame.add_paragraph()
    p.text = f"Answer: {mcq['answer']}"
    set_font_properties(p, 24, RGBColor(0, 102, 204))

    # Check the shape type
    # if answer_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Answer Box is a TextBox.")
    # else:
    #     print(f"The shape of Answer Box is not a TextBox, it is of type {answer_box.shape_type}")

# Function to add a background image to a slide
def add_background_image(slide, image_path, prs):
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)

# Function to add the first slide
def add_first_slide(prs, background_image_path, title_text, subtitle_text):
    first_slide_layout = prs.slide_layouts[6]  # Use a blank slide layout
    first_slide = prs.slides.add_slide(first_slide_layout)
    add_background_image(first_slide, background_image_path, prs)

    # Adding title
    title_shape = first_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    title_text_frame = title_shape.text_frame
    title_text_frame.text = title_text
    title_text_frame.paragraphs[0].font.size = Pt(56)
    title_text_frame.paragraphs[0].font.bold = True
    title_text_frame.paragraphs[0].font.color.rgb = RGBColor(20, 150, 25)
    title_text_frame.word_wrap = True  # Enable word wrap

    # Adding subtitle
    subtitle_shape = first_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    subtitle_text_frame = subtitle_shape.text_frame
    subtitle_text_frame.text = subtitle_text
    subtitle_text_frame.paragraphs[0].font.size = Pt(30)
    subtitle_text_frame.paragraphs[0].font.color.rgb = RGBColor(55, 125, 55)
    subtitle_text_frame.word_wrap = True  # Enable word wrap

    # Move the first slide to the beginning
    xml_slides = prs.slides._sldIdLst  # Access the slide ID list
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(0, slides[-1])

# Function to add the last slide
def add_last_slide(prs, background_image_path, last_slide_message):
    last_slide_layout = prs.slide_layouts[6]  # Use a blank slide layout
    last_slide = prs.slides.add_slide(last_slide_layout)
    add_background_image(last_slide, background_image_path, prs)

    # Adding message
    message_shape = last_slide.shapes.add_textbox(Inches(4), Inches(5), Inches(9), Inches(3))
    message_text_frame = message_shape.text_frame
    message_text_frame.text = last_slide_message
    message_text_frame.paragraphs[0].font.size = Pt(40)
    message_text_frame.paragraphs[0].font.bold = True
    message_text_frame.paragraphs[0].font.color.rgb = RGBColor(20, 150, 25)
    message_text_frame.word_wrap = True  # Enable word wrap

# Function to generate the PowerPoint with MCQs, first and last slides
def generate_ppt_with_first_last_slides(config):
    prs = Presentation()

    # Set the slide orientation to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add first slide at the beginning
    add_first_slide(prs, config['first_slide_bg'], config['first_slide_title_text'], config['first_slide_subtitle_text'])

    # Parse MCQs from the text file
    mcqs = parse_mcqs(config['mcq_file_path'])

    # Add slides with MCQs
    for index, mcq in enumerate(mcqs):
        add_slide(prs, mcq, index, config['all_slides_title'], config['watermark_path'])

    # Add last slide at the end
    add_last_slide(prs, config['last_slide_bg'], config['last_slide_message'])

    # Save the PowerPoint presentation
    prs.save(config['ppt_output_path_v1'])
    print(f"PowerPoint presentation saved as {config['ppt_output_path_v1']}")

def main():
    # Load configuration
    config_file = 'config.json'
    config = load_config(config_file)
    if config is None:
        return  # Exit if config could not be loaded

    # Generate PowerPoint presentation with the loaded configuration
    generate_ppt_with_first_last_slides(config)

if __name__ == "__main__":
    main()
