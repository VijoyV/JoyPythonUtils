import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Function to load the configuration from a JSON file
def load_config(config_file):
    try:
        with open(config_file, 'r', encoding='utf-8') as file:
            config = json.load(file)
        return config
    except FileNotFoundError:
        print(f"Error: Config file not found at {config_file}")
        return None
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}")
        return None

# Function to parse MCQs from a text file
def parse_mcqs(file_path):
    mcqs = []
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            lines = file.readlines()
    except FileNotFoundError:
        print(f"Error: File not found at {file_path}")
        return mcqs
    except Exception as e:
        print(f"Error reading file: {e}")
        return mcqs

    mcq = {}
    for line in lines:
        line = line.strip()
        if line.startswith("Question:"):
            if mcq:
                mcqs.append(mcq)
                mcq = {}
            mcq['question'] = line.split(':', 1)[1].strip()
        elif line.startswith("Options:"):
            mcq['options'] = []
        elif line.startswith('A)') or line.startswith('B)') or line.startswith('C)') or line.startswith('D)'):
            mcq['options'].append(line)
        elif line.startswith("Answer:"):
            mcq['answer'] = line.split(':', 1)[1].strip()
    if mcq:
        mcqs.append(mcq)
    return mcqs

# Function to set font properties for text paragraphs
def set_font_properties(paragraph, size, color=None):
    paragraph.font.size = Pt(size)
    if color:
        paragraph.font.color.rgb = color

# Function to add a slide with an MCQ to the presentation
def add_slide(prs, mcq, index, slide_title, watermark_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Add watermark image
    slide.shapes.add_picture(watermark_path, Inches(0), Inches(0),
                             width=prs.slide_width, height=prs.slide_height)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.1), Inches(10), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title = title_frame.add_paragraph()
    title.text = slide_title
    set_font_properties(title, 30, RGBColor(102, 55, 104))

    # Check the shape type
    # if title_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Title Box is a TextBox.")
    # else:
    #     print(f"The shape of Title Box is not a TextBox, it is of type {title_box.shape_type}")


    # Add question
    question_box = slide.shapes.add_textbox(Inches(1.1), Inches(1), Inches(10), Inches(2))
    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    p = question_frame.add_paragraph()
    p.text = f"Question - {index + 1}: {mcq['question']}"
    set_font_properties(p, 28, RGBColor(10, 10, 22))

    # Check the shape type
    # if question_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Question Box is a TextBox.")
    # else:
    #     print(f"The shape of Question Box is not a TextBox, it is of type {question_box.shape_type}")


    # Add options
    options_box = slide.shapes.add_textbox(Inches(1.1), Inches(3.5), Inches(10), Inches(3))
    options_frame = options_box.text_frame
    options_frame.word_wrap = True
    for option in mcq['options']:
        p = options_frame.add_paragraph()
        p.text = option
        p.level = 0
        set_font_properties(p, 28)

    # Check the shape type
    # if options_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Options Box is a TextBox.")
    # else:
    #     print(f"The shape of Options Box is not a TextBox, it is of type {options_box.shape_type}")

    # Add answer
    answer_box = slide.shapes.add_textbox(Inches(1.1), Inches(6), Inches(10), Inches(1))
    answer_frame = answer_box.text_frame
    answer_frame.word_wrap = True
    p = answer_frame.add_paragraph()
    p.text = f"Answer: {mcq['answer']}"
    set_font_properties(p, 28, RGBColor(4, 4, 4))

    # Check the shape type
    # if answer_box.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
    #     print("The shape of Answer Box is a TextBox.")
    # else:
    #     print(f"The shape of Answer Box is not a TextBox, it is of type {answer_box.shape_type}")

# Function to add a background image to a slide
def add_background_image(slide, image_path, prs):
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)

# Function to generate the PowerPoint with MCQs, first and last slides
def generate_ppt_with_mcqs(config):
    prs = Presentation()

    # Set the slide orientation to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Parse MCQs from the text file
    mcqs = parse_mcqs(config['mcq_file_path'])

    # Add slides with MCQs
    for index, mcq in enumerate(mcqs):
        add_slide(prs, mcq, index, config['all_slides_title'], config['watermark_path'])

    # Save the PowerPoint presentation
    prs.save(config['ppt_output_path_v1'])
    print(f"PowerPoint presentation saved as {config['ppt_output_path_v1']}")

def main():
    # Load configuration
    config_file = 'config.json'
    config = load_config(config_file)
    if config is None:
        print('Exits as config.json not found..!')
        return  # Exit if config could not be loaded

    # Generate PowerPoint presentation with the loaded configuration
    generate_ppt_with_mcqs(config)

if __name__ == "__main__":
    main()
