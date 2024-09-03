import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_shape_type(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "AutoShape"
    elif shape.shape_type == MSO_SHAPE_TYPE.CALLOUT:
        return "Callout"
    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return "Chart"
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "Picture"
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return "Table"
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "Group"
    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "Placeholder"
    elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
        return "Media"
    elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
        return "Line"
    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return "Text Box"
    else:
        return "Other"


def print_shapes_in_presentation(pptx_path):
    pptx_path = os.path.abspath(pptx_path)

    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"The file {pptx_path} does not exist.")

    try:
        presentation = Presentation(pptx_path)
        print("Presentation opened successfully.")

        for slide_number, slide in enumerate(presentation.slides, start=1):
            print(f"\nSlide {slide_number}:")

            for shape_number, shape in enumerate(slide.shapes, start=1):
                try:
                    shape_type = get_shape_type(shape)
                    print(f"  Shape {shape_number} - Type: {shape_type}")

                    # Check if the shape can have text and print it if it does
                    if shape.has_text_frame and shape.text.strip():
                        print(f"    Text: {shape.text.strip()}")
                except Exception as e:
                    print(f"  Error processing shape {shape_number}: {e}")

    except Exception as e:
        print(f"An error occurred: {e}")


if __name__ == "__main__":
    pptx_path = "Judges_Chapter_04_MCQ_V2.pptx"  # Replace with your PPTX file path
    print_shapes_in_presentation(pptx_path)
